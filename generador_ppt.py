import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ---
# 함수: TTF 파일에서 폰트 이름 추출

# ---
# 파워포인트 생성 함수 (한국어만)
# ---
def crear_ppt(titulos_kr, bloques_dict, secuencia, estilos, resaltados):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, titulo in enumerate(titulos_kr):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # ✅ slide 새로 만듦
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*estilos['bg_letra'])
        
        tb = slide.shapes.add_textbox(Inches(1), Inches(estilos['altura_texto']), Inches(11.33), Inches(3))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        run1 = p1.add_run()
        run1.text = titulo
        run1.font.size = Pt(estilos['tamano_titulo_kr'])
        run1.font.color.rgb = RGBColor(*estilos['color_titulo_kr'])
        p1.alignment = PP_ALIGN.CENTER

        for bloque_id in secuencia[i]:
            lineas = bloques_dict[i].get(bloque_id, [])
            for linea in lineas:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # ✅ slide 새로 만듦
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(*estilos['bg_letra'])
                
                tb = slide.shapes.add_textbox(Inches(1), Inches(estilos['altura_texto']), Inches(11.33), Inches(3))
                tf = tb.text_frame
                tf.clear()
                tf.word_wrap = True

                p1 = tf.paragraphs[0]
                run1 = p1.add_run()
                run1.text = linea
                run1.font.size = Pt(estilos['tamano_letra_kr'])
            
                # 💡 Aquí se aplica el color especial si es bloque resaltado
                if bloque_id == resaltados[i] and bloque_id != "":
                    run1.font.color.rgb = RGBColor(255, 192, 0)  # Dorado #FFC000
                else:
                    run1.font.color.rgb = RGBColor(*estilos['color_letra_kr'])

                p1.alignment = PP_ALIGN.CENTER

    return prs


# --- Streamlit UI ---
st.set_page_config(layout="wide")
st.title("마한장 (블록 반복 버전)")

num_canciones = st.number_input("찬양 개수", min_value=1, max_value=10, step=1)
altura_texto = st.slider("글자 위치 (0.0이 제일 높음)", 0.0, 6.0, value=1.0, step=0.1)

color_titulo_kr = "#000000"
bg_titulo = "#FFFFFF"
color_letra_kr = "#FFFFFF"
bg_letra = "#000000"

size_titulo_kr = st.number_input("[제목] 한국어 글자 크기", value=36)
size_letra_kr = st.number_input("[가사]  한국어 글자 크기", value=36)

estilos = {
    'color_titulo_kr': tuple(int(color_titulo_kr[i:i+2], 16) for i in (1, 3, 5)),
    'bg_titulo': tuple(int(bg_titulo[i:i+2], 16) for i in (1, 3, 5)),
    'altura_texto': altura_texto,

    'color_letra_kr': tuple(int(color_letra_kr[i:i+2], 16) for i in (1, 3, 5)),
    'bg_letra': tuple(int(bg_letra[i:i+2], 16) for i in (1, 3, 5)),
    'altura_texto': altura_texto,
    'tamano_titulo_kr': size_titulo_kr,
    'tamano_letra_kr': size_letra_kr,
}

korean_titles, bloques_por_cancion, secuencias, resaltados = [], [], [], []

for i in range(num_canciones):
    st.subheader(f"🎵 찬양 {i+1}")
    titulo = st.text_input(f"한국어 [제목] #{i+1}", key=f"kr_title_{i}")
    korean_titles.append(titulo)

    num_bloques = st.number_input(f"블록 수 #{i+1}", min_value=1, max_value=10, value=3, key=f"num_bloques_{i}")
    bloques = {}
    for j in range(num_bloques):
        nombre_bloque = st.text_input(f"🔠 블록 이름 #{j+1}", key=f"bloque_nombre_{i}_{j}")
        st.markdown(f"**✏️ {nombre_bloque} 가사**")
        contenido = st.text_area("", key=f"bloque_contenido_{i}_{j}")
        bloques[nombre_bloque] = contenido.split("\n")
    bloques_por_cancion.append(bloques)

    secuencia_str = st.text_input(f"슬라이드 순서 (예: A,A,B,C)", key=f"secuencia_{i}")
    bloque_resaltado = st.text_input(f"🎨 강조할 블록 이름 (선택사항)", key=f"resaltado_{i}").strip()
    resaltados.append(bloque_resaltado)
    secuencia = [s.strip() for s in secuencia_str.split(",") if s.strip() in bloques]
    secuencias.append(secuencia)

if st.button("🎷 PPT 생성"):
    it_path = il_path = None

    ppt = crear_ppt(korean_titles, bloques_por_cancion, secuencias, estilos, resaltados)

    ppt_path = "ppt_generado.pptx"
    ppt.save(ppt_path)

    with open(ppt_path, "rb") as f:
        st.download_button("📥 PPT 다운로드", f, file_name=ppt_path)
        
    if os.path.exists(ppt_path):
        os.remove(ppt_path)
